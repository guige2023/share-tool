#!/usr/bin/env python3
"""Office document text extractor: docx, xlsx, pptx → plain text JSON"""
import sys, json, os

def extract_docx(path):
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return {'text': '\n'.join(paragraphs), 'slides': len(paragraphs)}
    except Exception as e:
        return {'error': str(e)}

def extract_xlsx(path):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        sheets = []
        for name in wb.sheetnames[:5]:  # limit to 5 sheets
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(max_row=100, values_only=True):
                line = ' | '.join(str(c) if c is not None else '' for c in row)
                if line.strip():
                    rows.append(line)
            if rows:
                sheets.append({'name': name, 'rows': rows[:50]})
        text = '\n'.join(f"[{s['name']}]\n" + '\n'.join(s['rows']) for s in sheets)
        return {'text': text, 'sheets': len(sheets)}
    except Exception as e:
        return {'error': str(e)}

def extract_pptx(path):
    try:
        import pptx
        prs = pptx.Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides[:20]):  # limit to 20 slides
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append('\n'.join(texts))
        return {'text': '\n---\n'.join(slides), 'slides': len(slides)}
    except Exception as e:
        return {'error': str(e)}

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(json.dumps({'error': 'Usage: office-preview.py <docx|xlsx|pptx> <path>'}))
        sys.exit(1)

    fmt = sys.argv[1].lower()
    path = sys.argv[2]

    if not os.path.exists(path):
        print(json.dumps({'error': 'File not found'}))
        sys.exit(1)

    if fmt == 'docx':
        result = extract_docx(path)
    elif fmt == 'xlsx':
        result = extract_xlsx(path)
    elif fmt == 'pptx':
        result = extract_pptx(path)
    else:
        result = {'error': f'Unsupported format: {fmt}'}

    print(json.dumps(result, ensure_ascii=False))
